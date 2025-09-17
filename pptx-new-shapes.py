from pptx import Presentation
from pptx.util import Inches, Pt  # Inches用于设置位置和大小
from pptx.enum.shapes import MSO_SHAPE  # 形状类型
from pptx.dml.color import RGBColor  # 用于设置RGB颜色
from pptx.enum.text import PP_ALIGN  # 用于设置对齐方式

prs = Presentation()
slide_layout = prs.slide_layouts[5]  # 空白版式
slide = prs.slides.add_slide(slide_layout)

# 添加矩形（左、上、宽、高，单位：英寸）
left = Inches(1)
top = Inches(1)
width = Inches(3)
height = Inches(2)
rect = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,  # 矩形
    left, top, width, height
)

# 矩形填充颜色
rect.fill.solid()
rect.fill.fore_color.rgb = RGBColor(255, 255, 200)  # 浅黄色

# 矩形边框
rect.line.color.rgb = RGBColor(0, 0, 0)  # 黑色边框
rect.line.width = Pt(2)  # 边框粗细

# 形状内添加文本
tf = rect.text_frame
tf.text = "这是一个矩形"
tf.paragraphs[0].font.size = Pt(14)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER  # 文本居中

# 添加圆形
circle = slide.shapes.add_shape(
    MSO_SHAPE.OVAL,  # 圆形（椭圆的特殊形式）
    Inches(5), Inches(1), Inches(3), Inches(3)
)
circle.fill.solid()
circle.fill.fore_color.rgb = RGBColor(200, 255, 200)  # 浅绿色

prs.save("shapes.pptx")
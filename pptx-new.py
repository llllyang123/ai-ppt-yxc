from pptx import Presentation

# 创建空白演示文稿
prs = Presentation()

# 幻灯片版式：0-标题页，1-标题+内容，2-节标题，3-两栏内容，4-仅标题，5-空白等
# 添加标题页（版式0）
slide_layout = prs.slide_layouts[0]  # 选择版式
slide = prs.slides.add_slide(slide_layout)  # 添加幻灯片

# 获取标题和副标题占位符（版式0包含两个占位符）
title = slide.shapes.title
subtitle = slide.placeholders[1]  # 索引1对应副标题

# 设置文本
title.text = "python-pptx 教程"
subtitle.text = "—— 用Python自动生成PPT"

# 保存文件
prs.save("first_ppt.pptx")
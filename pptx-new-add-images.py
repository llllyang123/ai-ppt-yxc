from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)

# 插入本地图片（需提供图片路径）
img_path = "example.jpg"  # 替换为你的图片路径
left = Inches(1)
top = Inches(1)
height = Inches(5)  # 固定高度，宽度按比例缩放
pic = slide.shapes.add_picture(img_path, left, top, height=height)

# 添加图片标题
title = slide.shapes.title
title.text = "插入图片示例"

prs.save("image_slide.pptx")
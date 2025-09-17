from pptx import Presentation
from pptx.util import Pt  # 用于设置字体大小
from pptx.dml.color import RGBColor  # 用于设置RGB颜色
from pptx.enum.text import PP_ALIGN  # 用于设置对齐方式

prs = Presentation()
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

# 设置标题样式
title.text = "文本样式示例"
title_font = title.text_frame.paragraphs[0].font
title_font.name = "微软雅黑"  # 字体
title_font.size = Pt(32)  # 大小（磅）
title_font.color.rgb = RGBColor(0, 0, 128)  # 深蓝色（RGB值）
title_font.bold = True  # 加粗

# 设置内容样式
content.text_frame.text = "基础文本"
# 第一段（已通过text设置）
p1 = content.text_frame.paragraphs[0]
p1.font.name = "宋体"
p1.font.size = Pt(14)
p1.font.color.rgb = RGBColor(0, 0, 0)  # 黑色

# 第二段
p2 = content.text_frame.add_paragraph()
p2.text = "居中对齐 + 红色"
p2.alignment = PP_ALIGN.CENTER  # 居中对齐
p2.font.color.rgb = RGBColor(255, 0, 0)  # 红色

# 第三段
p3 = content.text_frame.add_paragraph()
p3.text = "右对齐 + 斜体"
p3.alignment = PP_ALIGN.RIGHT  # 右对齐
p3.font.italic = True  # 斜体

prs.save("text_style.pptx")
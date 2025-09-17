import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData  # 导入图表数据类


def json_to_ppt(json_data, output_path = "output.pptx"):
	"""
	将JSON数据转换为PPT文件（修复图表系列值设置错误）
	:param json_data: 包含PPT结构的JSON字典
	:param output_path: 输出PPT路径
	"""
	prs = Presentation()
	
	for slide_info in json_data["slides"]:
		layout_idx = slide_info["layout"]
		slide_layout = prs.slide_layouts[layout_idx]
		slide = prs.slides.add_slide(slide_layout)
		
		# 设置标题
		if "title" in slide_info:
			if slide.shapes.title:
				title = slide.shapes.title
			else:
				# 空白版式手动添加标题
				left = Inches(1)
				top = Inches(0.5)
				width = Inches(8)
				height = Inches(1)
				title_box = slide.shapes.add_textbox(left, top, width, height)
				title = title_box.text_frame
				title.paragraphs[0].alignment = PP_ALIGN.CENTER
			
			title.text = slide_info["title"]
			# 设置标题样式
			if hasattr(title, "text_frame"):
				para = title.text_frame.paragraphs[0]
			else:
				para = title.paragraphs[0]
			para.font.name = "微软雅黑"
			para.font.size = Pt(32)
			para.font.bold = True
		
		# 处理副标题（标题页版式）
		if layout_idx == 0 and "subtitle" in slide_info:
			subtitle = slide.placeholders[1]
			subtitle.text = slide_info["subtitle"]
			subtitle.text_frame.paragraphs[0].font.size = Pt(18)
		
		# 处理文本内容（标题+内容版式）
		if layout_idx == 1 and "content" in slide_info:
			content_placeholder = slide.placeholders[1]
			tf = content_placeholder.text_frame
			tf.clear()
			
			for idx, item in enumerate(slide_info["content"]):
				if idx == 0:
					p = tf.paragraphs[0]
				else:
					p = tf.add_paragraph()
				p.text = item
				p.font.name = "宋体"
				p.font.size = Pt(14)
				p.level = 0
		
		# 处理图表（修复部分）
		if "chart" in slide_info:
			chart_data = slide_info["chart"]
			left = Inches(1)
			top = Inches(2)
			width = Inches(8)
			height = Inches(5)
			
			# 创建图表数据对象
			ppt_chart_data = ChartData()
			ppt_chart_data.categories = chart_data["categories"]
			
			# 添加数据系列（修复值设置方式）
			for series in chart_data["series"]:
				# 正确方式：在添加系列时直接指定值
				ppt_chart_data.add_series(series["name"], series["values"])
			
			# 插入图表
			if chart_data["type"] == "column":
				chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
			elif chart_data["type"] == "line":
				chart_type = XL_CHART_TYPE.LINE
			else:
				chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
			
			chart = slide.shapes.add_chart(
				chart_type, left, top, width, height, ppt_chart_data
			).chart
			chart.has_title = True
			chart.chart_title.text_frame.text = "数据趋势图"
		
		# 处理图片
		if "image" in slide_info:
			img_path = slide_info["image"]
			left = Inches(1)
			top = Inches(2)
			height = Inches(5)
			try:
				slide.shapes.add_picture(img_path, left, top, height = height)
			except Exception as e:
				print(f"图片插入失败：{e}")
	
	prs.save(output_path)
	print(f"PPT已生成：{output_path}")


if __name__ == "__main__":
	sample_json = {
		"title": "AI生成的演示文稿",
		"slides": [
			{
				"layout": 0,
				"title": "AI技术发展报告",
				"subtitle": "2025年Q1"
			},
			{
				"layout": 1,
				"title": "核心技术方向",
				"content": [
					"大语言模型（LLM）的多模态融合",
					"AI Agents的自主决策能力提升",
					"边缘计算与AI的协同优化"
				]
			},
			{
				"layout": 5,
				"title": "市场规模预测（单位：亿元）",
				"chart": {
					"type": "column",
					"categories": ["2023", "2024", "2025", "2026"],
					"series": [
						{"name": "全球市场", "values": [800, 1200, 1800, 2500]},
						{"name": "中国市场", "values": [300, 500, 800, 1200]}
					]
				}
			}
		]
	}
	
	json_to_ppt(sample_json, "ai_generated_ppt.pptx")

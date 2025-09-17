from pptx import Presentation

prs = Presentation()

# 添加封面页（同上）
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Python-pptx 核心功能"
slide.placeholders[1].text = "基础操作指南"

# 添加内容页（版式1：标题+内容）
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "1. 文本处理"  # 标题

# 内容占位符（索引1）
content_placeholder = slide.placeholders[1]
tf = content_placeholder.text_frame  # 获取文本框

# 添加要点（段落）
tf.text = "• 支持设置字体、大小、颜色"  # 第一点
p = tf.add_paragraph()  # 新增段落
p.text = "• 支持对齐方式（左对齐、居中、右对齐）"
p = tf.add_paragraph()
p.text = "• 支持项目符号和编号"

prs.save("content_slide.pptx")